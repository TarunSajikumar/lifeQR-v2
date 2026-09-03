"""
LifeQR Master Real-World Ecosystem & Architecture PowerPoint Generator
Generates the definitive 20-slide, 16:9 widescreen presentation deck using python-pptx.
Merges ALL user requirements, objectives, hardware devices, process architectures,
and high-resolution diagrams:
- Stakeholder Matrix (Patient, Paramedic Crew, Hospital Doctor, Clinical Receptionist)
- 1-Tap SOS with Proximity Ambulance Dispatch
- Universal Hospital Passport (Free Gov & Private Clinic Interoperability)
- Unconscious / Unrecognizable Biker Crash Triage
- Dedicated Paramedic Tactical Hardware Device
- Doctor Desktop USB/BLE Scanner Cradle
- Doctor-to-Doctor Transfer & Specialist Referral Network
- Hospital Receptionist Advance Validation (Doctor Availability & Trauma Bay 1 Staging)
- School Student Chronic Illness Protection
- NTAG 424 DNA Anti-Cloning & Offline EEPROM Decryption
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_master_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    COLOR_BG_DARK = RGBColor(11, 14, 19)       # #0b0e13
    COLOR_CARD_DARK = RGBColor(20, 24, 33)     # #141821
    COLOR_BORDER = RGBColor(45, 52, 66)        # #2d3442
    COLOR_WHITE = RGBColor(255, 255, 255)      # #ffffff
    COLOR_RED = RGBColor(225, 29, 46)          # #E11D2E
    COLOR_CYAN = RGBColor(14, 165, 233)        # #0ea5e9
    COLOR_GREEN = RGBColor(16, 185, 129)       # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)       # #f59e0b
    COLOR_GRAY_LIGHT = RGBColor(209, 213, 219) # #d1d5db
    COLOR_GRAY_MUTED = RGBColor(156, 163, 175) # #9ca3af

    assets_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/assets'))

    def set_slide_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        pill_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.7), Inches(0.32))
        tf_pill = pill_box.text_frame
        tf_pill.word_wrap = True
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = f"●  LIFEQR MASTER ECOSYSTEM  |  {category.upper()}"
        p_pill.font.size = Pt(9.5)
        p_pill.font.bold = True
        p_pill.font.color.rgb = COLOR_RED

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        if subtitle:
            p_sub = tf_title.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(11.5)
            p_sub.font.color.rgb = COLOR_GRAY_MUTED

    # =========================================================================
    # SLIDE 1: Cover Slide
    # =========================================================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s1, COLOR_BG_DARK)

    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.18), Inches(5.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_RED
    bar.line.fill.background()

    b1 = s1.shapes.add_textbox(Inches(1.25), Inches(1.1), Inches(11.0), Inches(5.2))
    tf1 = b1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "LIFEQR: REAL-WORLD CONNECTED EMERGENCY ECOSYSTEM"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Saving Lives as Quick as Possible: Unified Patient, Paramedic, and Hospital Doctor Network"
    p2.font.size = Pt(17)
    p2.font.color.rgb = COLOR_RED
    p2.space_before = Pt(8)

    p3 = tf1.add_paragraph()
    p3.text = "An end-to-end physical-to-digital lifesaver:\n• Patient: 1-Tap SOS, Universal Medical Passport (Gov & Private), Hybrid NFC/QR Smart Card\n• Paramedic Crew: Tactical Handheld Terminal, Unrecognizable Biker Triage, Offline Vitals Decryption\n• Hospital & Doctor: Advance Arrival Telemetry, Receptionist Validation, Desktop Cradle & Doctor Transfers"
    p3.font.size = Pt(12.5)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.space_before = Pt(18)

    p4 = tf1.add_paragraph()
    p4.text = "EXECUTIVE PITCH  •  FULL-STACK PROCESS ARCHITECTURE  •  HARDWARE SPECIFICATIONS"
    p4.font.size = Pt(10.5)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GRAY_MUTED
    p4.space_before = Pt(40)

    # =========================================================================
    # SLIDE 2: Core Mission & The Real-World Crisis
    # =========================================================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s2, COLOR_BG_DARK)
    add_header(s2, "Core Mission", "The Real-World Crisis: The Golden Hour Dilemma", "Why speed is the single deciding factor between survival and preventable death")

    cards2 = [
        ("The 'Golden Hour' Reality", "Trauma mortality spikes by over 300% if surgical resuscitation is delayed beyond the initial 60 minutes after severe impact.", "CRITICAL WINDOW", COLOR_RED),
        ("70% Unconscious Victims", "Accident victims, seizing students, or stroke patients cannot speak, communicate allergies, or unlock their phones.", "PATIENT VOID", COLOR_AMBER),
        ("Unrecognizable Biker Accidents", "High-speed highway crashes leave riders unconscious with heavy gear on; removing helmets risks fatal spinal severance.", "BIKER DANGER", COLOR_RED),
        ("Disconnected ER Desks", "Ambulances rush to clinics only to discover no open trauma bay or on-duty surgeon is ready, wasting vital minutes.", "HOSPITAL DELAY", COLOR_AMBER)
    ]

    for i, (head, body, tag, col) in enumerate(cards2):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = head
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 3: The 3 Core Stakeholders
    # =========================================================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s3, COLOR_BG_DARK)
    add_header(s3, "Stakeholders", "The 3 Stakeholder Roles in LifeQR", "Empowering patients, paramedics, and doctors in a continuous lifesaver loop")

    stakeholders = [
        ("1. CITIZEN & PATIENT", "Universal Medical Passport", "• 1-Tap SOS sends precise GPS to nearby ambulances.\n• Hybrid NFC/QR card carried in wallet.\n• Helmet decals for riders; bag tags for children.\n• Accepted across any Government or Private hospital with zero clinical confusion.", COLOR_GREEN),
        ("2. PARAMEDIC & AMBULANCE", "Tactical Triage Crew", "• Rugged handheld device for driver & nurse.\n• Instant 0.2s dual NFC/QR scan of unconscious victims.\n• Offline on-card EEPROM decryption of blood group & allergies.\n• Live GPS hospital detection & en route reception validation.", COLOR_RED),
        ("3. DOCTOR & HOSPITAL", "Clinical Reception & Care", "• Desktop USB/BLE scanner cradle for 1-tap intake.\n• Advance notification desk: validates doctor availability.\n• Trauma Bay 1 lockout & blood bank staging.\n• Seamless doctor-to-doctor transfer and referral channel.", COLOR_CYAN)
    ]

    for i, (badge, title, desc, col) in enumerate(stakeholders):
        x = Inches(0.8 + (i * 3.95))
        y = Inches(2.0)
        c = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = badge
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 4: Master System Architecture (With Image)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s4, COLOR_BG_DARK)
    add_header(s4, "System Topology", "Full-Stack Connected Emergency Ecosystem", "The physical-to-cloud infrastructure uniting citizen, responder, and trauma center")

    img_top = os.path.join(assets_dir, 'lifeqr_system_architecture.jpg')
    if os.path.exists(img_top):
        s4.shapes.add_picture(img_top, Inches(0.8), Inches(1.85), width=Inches(7.2), height=Inches(5.1))

    c4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(1.85), Inches(4.3), Inches(5.1))
    c4.fill.solid()
    c4.fill.fore_color.rgb = COLOR_CARD_DARK
    c4.line.color.rgb = COLOR_BORDER

    tf4 = c4.text_frame
    tf4.word_wrap = True
    tf4.margin_left = Inches(0.25)
    tf4.margin_top = Inches(0.25)
    tf4.margin_right = Inches(0.25)

    p = tf4.paragraphs[0]
    p.text = "THE 4-TIER ARCHITECTURE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    tiers_data = [
        ("Tier 1: Citizen Physical Edge", "NTAG 424 DNA smart cards, 3M helmet decals, and IMU crash beacons. Zero battery, permanent passive operation.", COLOR_GREEN),
        ("Tier 2: Tactical Responder Gateway", "Rugged ambulance terminal with dual NFC/2D optical engine, offline SQLite cache, 4G LTE-M & satellite fallback.", COLOR_RED),
        ("Tier 3: Cloud & Dispatch Core", "Microservices event bus (Socket.IO/Kafka), PostGIS spatial proximity routing, zero-knowledge AES-256 GCM vault.", COLOR_CYAN),
        ("Tier 4: Clinical Hospital ER", "HL7/FHIR v4 EHR interoperability, doctor desktop USB cradle, and live incoming trauma bay reservation radar.", COLOR_AMBER)
    ]

    for t_head, t_body, t_col in tiers_data:
        p_h = tf4.add_paragraph()
        p_h.text = f"■ {t_head}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = t_col
        p_h.space_before = Pt(8)

        p_b = tf4.add_paragraph()
        p_b.text = t_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = COLOR_GRAY_LIGHT
        p_b.space_before = Pt(2)

    # =========================================================================
    # SLIDE 5: Process 1 Blueprint — Patient Lifecycle (With Image)
    # =========================================================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s5, COLOR_BG_DARK)
    add_header(s5, "Process Architecture 01", "Patient & Citizen Lifecycle Blueprint", "Registration, NTAG 424 DNA chip encoding, laser QR printing, and zero-knowledge encryption")

    img_pat = os.path.join(assets_dir, 'patient_process_arch.jpg')
    if os.path.exists(img_pat):
        s5.shapes.add_picture(img_pat, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    c5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    c5.fill.solid()
    c5.fill.fore_color.rgb = COLOR_CARD_DARK
    c5.line.color.rgb = COLOR_GREEN
    c5.line.width = Pt(1.5)

    tf5 = c5.text_frame
    tf5.word_wrap = True
    tf5.margin_left = Inches(0.2)
    tf5.margin_top = Inches(0.25)
    tf5.margin_right = Inches(0.2)

    p = tf5.paragraphs[0]
    p.text = "PROCESS 1 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    p_highlights = [
        ("Step 1: Intake", "Blood group, allergies, chronic ailments, and primary ICE contacts recorded."),
        ("Step 2: NTAG 424 DNA", "Card issuance encoder writes rolling AES-128 SUN keys to prevent counterfeiting."),
        ("Step 3: Wearables", "Wallet card, helmet decal, and student backpack tag deployed."),
        ("Step 4: Vault", "Data secured via AES-256 GCM client-side encryption.")
    ]

    for h_t, h_d in p_highlights:
        p_h = tf5.add_paragraph()
        p_h.text = f"✔ {h_t}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf5.add_paragraph()
        p_d.text = h_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 6: Universal Medical Passport (Gov & Private Hospitals)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s6, COLOR_BG_DARK)
    add_header(s6, "Universal Passport", "The Universal Hospital Medical Passport", "One physical identity card working seamlessly across any government or private clinic")

    cards6 = [
        ("Zero Hospital Walled Gardens", "Patients visit free public/government hospitals, private specialist centers, or urgent care clinics with identical zero-friction access.", "SEAMLESS MOBILITY", COLOR_GREEN),
        ("Instant Historical Medical Pre-View", "Doctors instantly view past surgery notes, chronic diseases, previous lab scans, and lifestyle alerts, eliminating diagnostic confusion.", "ELIMINATING CONFUSION", COLOR_CYAN),
        ("Zero Redundant Tests", "Prevents duplicate blood panels, redundant X-rays, or unnecessary MRI delays, dramatically reducing healthcare costs and treatment delays.", "SAVING EXPENSE & TIME", COLOR_AMBER),
        ("Patient-Controlled Consent", "Patient decides what data is shared; emergency profile is public, while full clinical consultation charts require patient authentication.", "PRIVACY BY DESIGN", COLOR_WHITE)
    ]

    for i, (head, body, tag, col) in enumerate(cards6):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = head
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 7: Wearables & Everyday Life Protection (With Image)
    # =========================================================================
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s7, COLOR_BG_DARK)
    add_header(s7, "Physical Wearables", "Wearables & Everyday Life Protection", "Protecting motorcyclists, school students, seniors, and everyday citizens")

    img_wear = os.path.join(assets_dir, 'nfc_card_wearables.jpg')
    if os.path.exists(img_wear):
        s7.shapes.add_picture(img_wear, Inches(0.8), Inches(1.9), width=Inches(6.6), height=Inches(4.9))

    c7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.9), Inches(4.9))
    c7.fill.solid()
    c7.fill.fore_color.rgb = COLOR_CARD_DARK
    c7.line.color.rgb = COLOR_GREEN
    c7.line.width = Pt(1.5)

    tf7 = c7.text_frame
    tf7.word_wrap = True
    tf7.margin_left = Inches(0.25)
    tf7.margin_top = Inches(0.25)
    tf7.margin_right = Inches(0.25)

    p = tf7.paragraphs[0]
    p.text = "PHYSICAL FORM FACTORS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    wear_items = [
        ("Wallet Smart Medical Card", "Waterproof composite PVC card with NTAG 424 DNA NFC chip + high-contrast laser QR code. Carried by everyday adults and commuters."),
        ("Motorcycle Helmet Crash Decal", "3M retroreflective vinyl sticker placed on rider helmets so bystanders or EMS can scan without removing helmets after spinal trauma."),
        ("School & Student Bag Tags", "Clip-on tag for school backpacks protecting children with epilepsy, asthma, diabetes, or nut allergies. Teachers scan in 1.8s if a child collapses."),
        ("Bystander 1-Tap Phone Access", "Any passerby taps with their phone $\rightarrow$ sees critical blood type and immediately clicks 'Call Family Contact (ICE)' with zero app install.")
    ]

    for w_h, w_d in wear_items:
        p_h = tf7.add_paragraph()
        p_h.text = f"✔ {w_h}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf7.add_paragraph()
        p_d.text = w_d
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 8: Patient 1-Tap SOS Dispatch Flow
    # =========================================================================
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s8, COLOR_BG_DARK)
    add_header(s8, "Emergency SOS", "Patient 1-Tap SOS & Ambulance Proximity Dispatch", "Instant automated dispatch broadcasting exact coordinates to on-duty ambulances")

    sos_steps = [
        ("1. 1-Tap SOS Trigger", "Patient or bystander clicks SOS on web app or scans NFC card. High-precision GPS coordinates locked in < 1 second.", COLOR_RED),
        ("2. Cloud Dispatch Engine", "LifeQR cluster queries on-duty ambulance units within a 15 km radius via WebSocket Socket.IO room `crew:all`.", COLOR_AMBER),
        ("3. Ambulance Terminal Alert", "Nearby ambulance driver's tactical device chimes loudly; screen displays live navigation route and victim triage level.", COLOR_CYAN),
        ("4. Automated ICE SMS Broadcast", "Emergency contacts (spouses, parents) receive automated SMS with live tracking link and ambulance call sign.", COLOR_GREEN)
    ]

    for i, (title, desc, col) in enumerate(sos_steps):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"STEP 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 9: Process 2 Blueprint — Bike Accident Triage (With Image)
    # =========================================================================
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s9, COLOR_BG_DARK)
    add_header(s9, "Process Architecture 02", "Highway Motorcycle Accident Triage Blueprint", "IMU crash detection, paramedic tactical dual-scan, offline EEPROM decode & GPS dispatch")

    img_acc = os.path.join(assets_dir, 'accident_process_arch.jpg')
    if os.path.exists(img_acc):
        s9.shapes.add_picture(img_acc, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    c9 = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    c9.fill.solid()
    c9.fill.fore_color.rgb = COLOR_CARD_DARK
    c9.line.color.rgb = COLOR_RED
    c9.line.width = Pt(1.5)

    tf9 = c9.text_frame
    tf9.word_wrap = True
    tf9.margin_left = Inches(0.2)
    tf9.margin_top = Inches(0.25)
    tf9.margin_right = Inches(0.2)

    p = tf9.paragraphs[0]
    p.text = "PROCESS 2 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    acc_highlights = [
        ("1. Motorcycle Crash", "IoT IMU sensor beacon detects >15G impact, broadcasting automated BLE emergency beacon."),
        ("2. Paramedic Arrives", "Arrives with rugged tactical handheld terminal; no personal smartphone delays."),
        ("3. Instant 0.2s Scan", "Scans helmet QR decal or taps wallet NFC card; decrypts O+ blood and allergy offline."),
        ("4. GPS CAD Dispatch", "Auto-identifies nearest Level-1 trauma center & transmits live vitals en route.")
    ]

    for a_t, a_d in acc_highlights:
        p_h = tf9.add_paragraph()
        p_h.text = f"★ {a_t}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf9.add_paragraph()
        p_d.text = a_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 10: Unrecognizable Biker Crash Triage
    # =========================================================================
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s10, COLOR_BG_DARK)
    add_header(s10, "Biker Triage", "Triage of Unrecognizable & Unconscious Bikers", "Overcoming face-obscuring helmets, shock, and missing physical identification")

    biker_points = [
        ("The Helmet Removal Risk", "Removing a full-face motorcycle helmet on an unconscious rider with cervical spine trauma can cause permanent paralysis or death. LifeQR helmet QR decal allows full triage scanning from the outside.", COLOR_RED),
        ("Unrecognizable Facial Identity", "Severe facial polytrauma makes biometric face unlock impossible. The NFC chip embedded in wallet or jacket sleeve resolves patient identity in 200 milliseconds.", COLOR_AMBER),
        ("Critical Contraindication Warnings", "Instantly flags lethal contraindications (e.g. administering standard anticoagulants to a hemophiliac or penicillin to an anaphylactic victim).", COLOR_GREEN),
        ("Paramedic Nurse Medication Lookup", "The ambulance nurse accesses the victim's verified daily medications in one touch, preventing fatal drug interactions during roadside resuscitation.", COLOR_CYAN)
    ]

    for i, (head, body, col) in enumerate(biker_points):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 11: Ambulance Tactical Terminal Hardware (With Image)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s11, COLOR_BG_DARK)
    add_header(s11, "Paramedic Hardware", "Ambulance Tactical Triage Terminal", "Dedicated field device engineered for ambulance drivers and paramedic nurses")

    img_dev = os.path.join(assets_dir, 'ambulance_triage_device.jpg')
    if os.path.exists(img_dev):
        s11.shapes.add_picture(img_dev, Inches(0.8), Inches(1.9), width=Inches(6.6), height=Inches(4.9))

    c11 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.9), Inches(4.9))
    c11.fill.solid()
    c11.fill.fore_color.rgb = COLOR_CARD_DARK
    c11.line.color.rgb = COLOR_RED
    c11.line.width = Pt(1.5)

    tf11 = c11.text_frame
    tf11.word_wrap = True
    tf11.margin_left = Inches(0.25)
    tf11.margin_top = Inches(0.25)
    tf11.margin_right = Inches(0.25)

    p = tf11.paragraphs[0]
    p.text = "TACTICAL FIELD ENGINE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    dev_features = [
        ("Dual Sensor Engine", "Top NFC RF halo + 60fps 2D CMOS barcode optical scanner reads cards, tags, and helmet stickers through rain or grime in < 0.2s."),
        ("Immediate Triage HUD", "Instantly projects Blood Group, Lethal Allergies, Current Medications, and Emergency Contacts on 800-nit high-contrast screen."),
        ("Automated Hospital Detection", "Integrated GNSS GPS queries hospital network, calculates traffic-adjusted ETA, and validates ER availability en route."),
        ("1-Touch ER Vitals Stream", "Nurse enters pulse, SpO2, and blood pressure with 1 touch, transmitting live vitals straight to the receiving trauma bay."),
        ("Rugged Industrial Design", "IP68 dust/waterproof, drop-resistant rubber bumper casing for harsh accident environments.")
    ]

    for d_h, d_d in dev_features:
        p_h = tf11.add_paragraph()
        p_h.text = f"★ {d_h}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf11.add_paragraph()
        p_d.text = d_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 12: Tactical Terminal Exploded Assembly & BOM (With Image)
    # =========================================================================
    s12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s12, COLOR_BG_DARK)
    add_header(s12, "Hardware Engineering", "Tactical Triage Terminal: Exploded Architecture", "Layer-by-layer industrial assembly design for ambulance field reliability")

    img_exp = os.path.join(assets_dir, 'hardware_exploded_view.jpg')
    if os.path.exists(img_exp):
        s12.shapes.add_picture(img_exp, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    c12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    c12.fill.solid()
    c12.fill.fore_color.rgb = COLOR_CARD_DARK
    c12.line.color.rgb = COLOR_AMBER
    c12.line.width = Pt(1.5)

    tf12 = c12.text_frame
    tf12.word_wrap = True
    tf12.margin_left = Inches(0.2)
    tf12.margin_top = Inches(0.25)
    tf12.margin_right = Inches(0.2)

    p = tf12.paragraphs[0]
    p.text = "ENGINEERING BOM"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    exp_specs = [
        ("[1] Armored Bezel", "High-impact composite protecting against face-down asphalt drops."),
        ("[2] Gorilla Glass 5", "Optical-bonded 800-nit sunlight-readable IPS LCD."),
        ("[3] Planar NFC Coil", "Tuned 13.56 MHz loop antenna for omnidirectional tap."),
        ("[4] 2D CMOS Engine", "Newland optical imager reading scratched QR codes."),
        ("[5] Multi-Layer PCB", "ESP32-S3 SoC, Quectel 4G LTE-M, GNSS GPS, ATECC608B."),
        ("[6] IP68 Chassis", "TPU rubber bumpers resisting hospital decontamination wipes.")
    ]

    for e_h, e_d in exp_specs:
        p_h = tf12.add_paragraph()
        p_h.text = f"• {e_h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(6)

        p_d = tf12.add_paragraph()
        p_d.text = e_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 13: Offline-First Operation & Satellite Fallback
    # =========================================================================
    s13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s13, COLOR_BG_DARK)
    add_header(s13, "Edge Resilience", "Offline-First Operation & Satellite Fallback", "Solving the highway dead zone: 100% triage uptime with zero cellular connectivity")

    off_points = [
        ("The Remote Highway Dead Zone", "Severe crashes often occur in rural mountain roads or tunnel corridors with 0% cellular signal. Cloud-only architectures fail catastrophically in these moments.", "REALITY OF CRASH SITES", COLOR_RED),
        ("Encrypted On-Card EEPROM Payload", "512 bytes of critical data (Blood Group, Top 3 Lethal Allergies, ICE Phone, Chronic Illness) stored directly inside the card's chip memory. Decrypted locally in < 0.1s with 0 internet.", "ON-CARD EDGE VAULT", COLOR_GREEN),
        ("Terminal Local SQLite Cache", "Ambulance terminal maintains an onboard encrypted local database of regional emergency protocols, drug interactions, and hospital directories updated daily via Wi-Fi.", "LOCAL INTELLIGENCE", COLOR_CYAN),
        ("Satellite & LoRa Mesh SOS Fallback", "When cellular fails, the terminal automatically falls back to LoRa emergency mesh or direct satellite burst (Iridium / Starlink Direct-to-Cell) to alert emergency dispatch.", "SATELLITE REDUNDANCY", COLOR_AMBER)
    ]

    for i, (head, body, tag, col) in enumerate(off_points):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s13.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = head
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 14: Process 3 Blueprint — Hospital ER Resuscitation (With Image)
    # =========================================================================
    s14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s14, COLOR_BG_DARK)
    add_header(s14, "Process Architecture 03", "Hospital & Clinical ER Resuscitation Blueprint", "In-transit telemetry, automated trauma bay reservation, doctor desktop cradle & FHIR sync")

    img_hosp = os.path.join(assets_dir, 'hospital_process_arch.jpg')
    if os.path.exists(img_hosp):
        s14.shapes.add_picture(img_hosp, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    c14 = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    c14.fill.solid()
    c14.fill.fore_color.rgb = COLOR_CARD_DARK
    c14.line.color.rgb = COLOR_CYAN
    c14.line.width = Pt(1.5)

    tf14 = c14.text_frame
    tf14.word_wrap = True
    tf14.margin_left = Inches(0.2)
    tf14.margin_top = Inches(0.25)
    tf14.margin_right = Inches(0.2)

    p = tf14.paragraphs[0]
    p.text = "PROCESS 3 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    hosp_highlights = [
        ("1. In-Transit Vitals", "ER monitor receives live countdown (ETA 07:12) and vitals (HR 134, SpO2 89%, BP 95/58)."),
        ("2. Bay 1 & Blood Staged", "Trauma Bay 1 locked; blood bank thaws matched O-negative units before siren arrives."),
        ("3. Doctor Desktop Cradle", "Doctor taps card on USB/BLE scanner cradle; profile populates with zero manual typing."),
        ("4. FHIR & AI Safety", "Auto-syncs into Epic/Cerner EHR; flags penicillin allergy and generates SOAP chart.")
    ]

    for h_t, h_d in hosp_highlights:
        p_h = tf14.add_paragraph()
        p_h.text = f"✔ {h_t}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf14.add_paragraph()
        p_d.text = h_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 15: Hospital Pre-Notification & Reception Validation
    # =========================================================================
    s15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s15, COLOR_BG_DARK)
    add_header(s15, "Reception Triage", "Hospital Pre-Notification & Reception Validation Flow", "Confirming doctor availability and clinical capacity before ambulance arrival")

    recept_steps = [
        ("1. Automated Hospital Siren Alert", "When ambulance locks destination, the hospital reception desk and ER triage nurse receive a flashing siren alert with victim summary.", COLOR_AMBER),
        ("2. Doctor Availability Validation", "Receptionist / ER Charge Nurse verifies specialized doctor availability (neurosurgeon, orthopedic trauma, cardiac surgeon on-duty).", COLOR_CYAN),
        ("3. Capacity Validation & Acceptance", "System verifies available ICU bed and Trauma Bay 1 readiness; transmits an authenticated confirmation token back to the ambulance.", COLOR_GREEN),
        ("4. Dynamic Diversion if at Capacity", "If ER is at max capacity (divert status), the system instantly alerts the paramedic and re-routes ambulance to the next nearest trauma center.", COLOR_RED)
    ]

    for i, (title, desc, col) in enumerate(recept_steps):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"STAGE 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 16: Trauma Bay 1 Reservation & Advance Blood Staging
    # =========================================================================
    s16 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s16, COLOR_BG_DARK)
    add_header(s16, "Trauma Resuscitation", "Trauma Bay 1 Lockout & Advance Blood Staging", "Zero-delay surgical handover prepared 7 minutes before the ambulance wheels stop")

    er_actions = [
        ("Trauma Bay Lockout (ETA -7m)", "Trauma Bay 1 is electronically locked on the hospital central board. Ventilator circuit, chest tube tray, and ultrasound FAST machine pre-staged at bedside.", COLOR_RED),
        ("Blood Bank Staging (ETA -5m)", "Automated reservation signal triggers the hospital blood bank refrigerator to thaw 2 units of uncrossed O-negative PRBCs and 2 units of FFP. Ready at bay bedside.", COLOR_CYAN),
        ("Trauma Team Summoned", "On-call trauma surgeon, anesthesiologist, respiratory therapist, and surgical techs alerted via priority pager with ETA countdown.", COLOR_AMBER),
        ("Immediate Dock Handover (ETA 0m)", "Ambulance backs into trauma bay; victim wheeled straight onto Bay 1 resuscitation gurney. Zero check-in registration delay. Resuscitation begins instantly.", COLOR_GREEN)
    ]

    for i, (title, desc, col) in enumerate(er_actions):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = s16.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 17: Doctor Desktop USB Scanner Cradle & 1-Tap Intake (With Image)
    # =========================================================================
    s17 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s17, COLOR_BG_DARK)
    add_header(s17, "Doctor Hardware", "Doctor Desktop Scanner & Zero-Typing Consultation Flow", "Eliminating transcription errors, medication mix-ups, and consultation paperwork delays")

    img_doc = os.path.join(assets_dir, 'doctor_desktop_scanner.jpg')
    if os.path.exists(img_doc):
        s17.shapes.add_picture(img_doc, Inches(0.8), Inches(1.9), width=Inches(6.6), height=Inches(4.9))

    c17 = s17.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.9), Inches(4.9))
    c17.fill.solid()
    c17.fill.fore_color.rgb = COLOR_CARD_DARK
    c17.line.color.rgb = COLOR_CYAN
    c17.line.width = Pt(1.5)

    tf17 = c17.text_frame
    tf17.word_wrap = True
    tf17.margin_left = Inches(0.25)
    tf17.margin_top = Inches(0.25)
    tf17.margin_right = Inches(0.25)

    p = tf17.paragraphs[0]
    p.text = "PHYSICIAN DESK WORKFLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    doc_points = [
        ("Plug-and-Play USB HID/CCID", "Scanner cradle connects via USB Type-C to doctor's PC. Functions driverless as a smart card reader and keyboard wedge."),
        ("1-Tap Patient Record Ingestion", "Doctor simply taps the patient's card or scans QR code $\rightarrow$ complete clinical profile, past surgeries, and active prescriptions populate with zero manual typing."),
        ("AI Clinical Safety Copilot", "Integrated AI cross-references newly prescribed medication against allergies and current drugs, instantly flagging lethal drug interactions."),
        ("Automated SOAP Note Generation", "Generates compliant Subjective, Objective, Assessment, Plan (SOAP) clinical encounter documentation ready for physician signature.")
    ]

    for d_h, d_d in doc_points:
        p_h = tf17.add_paragraph()
        p_h.text = f"✔ {d_h}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_d = tf17.add_paragraph()
        p_d.text = d_d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(2)

    # =========================================================================
    # SLIDE 18: Doctor-to-Doctor Transfer & Specialist Referral Network
    # =========================================================================
    s18 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s18, COLOR_BG_DARK)
    add_header(s18, "Doctor Collaboration", "Doctor-to-Doctor Transfer & Specialist Referral Network", "Connecting physicians across departments and hospitals with verified handover tokens")

    transfer_nodes = [
        ("1. In-App Patient Transfer", "Attending physician selects 'Refer / Transfer Patient' on the LifeQR Clinical Dashboard, choosing a specialized department or higher-level surgeon.", COLOR_CYAN),
        ("2. Clinical Handover Dossier", "System bundles active vitals, diagnostic notes, imaging reports, and current medication list into a secure cryptographic transfer token.", COLOR_GREEN),
        ("3. Receiving Specialist Notification", "Receiving surgeon receives an encrypted notification. Tapping patient's card or accepting the token imports the complete clinical history without re-interviewing the patient.", COLOR_AMBER),
        ("4. Inter-Doctor Chat & Tele-Consult", "Built-in encrypted physician-to-physician communication channel allows real-time consults, surgical second opinions, and ICU transfer coordination.", COLOR_WHITE)
    ]

    for i, (title, desc, col) in enumerate(transfer_nodes):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = s18.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"STEP 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 19: Enterprise Interoperability & Zero-Trust Security
    # =========================================================================
    s19 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s19, COLOR_BG_DARK)
    add_header(s19, "Security & Interoperability", "HL7/FHIR v4 Interoperability & Zero-Trust Security", "Universal hospital connectivity with HIPAA-compliant cryptographic governance")

    sec_blocks = [
        ("Standardized FHIR v4 Gateway", "Natively maps clinical data to standard HL7 FHIR resources (`Patient`, `AllergyIntolerance`, `Condition`, `Encounter`), immediately readable by Epic, Cerner, and OpenEMR without custom coding.", COLOR_CYAN),
        ("Role-Based Dynamic Masking", "• Tier A (Public Bystander): Blood group, ICE phone numbers, anaphylaxis warning.\n• Tier B (Paramedic on Scene): Active medications, chronic conditions, emergency directives.\n• Tier C (Doctor / ER): Complete clinical records, surgeries, lab tests.", COLOR_GREEN),
        ("Zero-Knowledge Vault", "Patient medical data is encrypted client-side using AES-256 GCM. Master keys are protected in hardware security modules (HSM) with zero plain-text leaks.", COLOR_AMBER),
        ("Cryptographic Audit Trail", "Every tap, scan, and clinical handover is cryptographically signed, timestamped, and geotagged. Immutable audit logs recorded in full compliance with HIPAA Section 164.312.", COLOR_RED)
    ]

    for i, (title, desc, col) in enumerate(sec_blocks):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = s19.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 20: Master Execution Roadmap & Conclusion
    # =========================================================================
    s20 = prs.slides.add_slide(blank_layout)
    set_slide_bg(s20, COLOR_BG_DARK)

    cta_card = s20.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.1), Inches(10.33), Inches(5.3))
    cta_card.fill.solid()
    cta_card.fill.fore_color.rgb = COLOR_CARD_DARK
    cta_card.line.color.rgb = COLOR_RED
    cta_card.line.width = Pt(2)

    tf20 = cta_card.text_frame
    tf20.word_wrap = True
    tf20.margin_left = Inches(0.6)
    tf20.margin_top = Inches(0.5)
    tf20.margin_right = Inches(0.6)

    p = tf20.paragraphs[0]
    p.text = "WHEN SECONDS COUNT, LIFEQR SAVES LIVES."
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf20.add_paragraph()
    p2.text = "A Complete, Synchronized, Real-World Emergency Medical Network"
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_RED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    p3 = tf20.add_paragraph()
    p3.text = "By uniting passive hybrid NFC/QR wearables, dedicated ambulance triage terminals,\nuniversal hospital medical passports (Gov & Private), doctor desktop scanner cradles,\nand automated ER pre-notification, LifeQR turns preventable tragedies into survivable emergencies.\nWe are ready to build the physical devices, distribute the wearables, and transform emergency healthcare."
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(18)

    p4 = tf20.add_paragraph()
    p4.text = "READY FOR HARDWARE FABRICATION & MUNICIPAL EMS INTEGRATION"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GREEN
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(26)

    # Save presentation
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/LifeQR_Master_RealWorld_Ecosystem.pptx'))
    prs.save(output_path)
    print(f"[Success] Master Real-World Ecosystem Presentation generated at: {output_path}")

if __name__ == '__main__':
    create_master_presentation()
