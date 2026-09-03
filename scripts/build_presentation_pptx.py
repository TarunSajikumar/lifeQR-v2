"""
LifeQR Real-World Ecosystem PowerPoint Generator
Generates a 12-slide, 16:9 widescreen presentation deck using python-pptx.
Incorporates real-world IoT hardware concepts, NFC wearables, ambulance triage,
doctor consultation scanners, and ER incoming hospital dispatch.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette (Swiss Brutalist / LifeQR High-Contrast)
    COLOR_BG_DARK = RGBColor(15, 17, 21)      # #0f1115
    COLOR_CARD_DARK = RGBColor(26, 29, 36)    # #1a1d24
    COLOR_BORDER = RGBColor(55, 65, 81)       # #374151
    COLOR_WHITE = RGBColor(255, 255, 255)     # #ffffff
    COLOR_RED = RGBColor(225, 29, 46)         # #E11D2E
    COLOR_GRAY_LIGHT = RGBColor(209, 213, 219)# #d1d5db
    COLOR_GRAY_MUTED = RGBColor(156, 163, 175)# #9ca3af
    COLOR_GREEN = RGBColor(16, 185, 129)      # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)      # #f59e0b

    assets_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/assets'))

    def set_slide_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        # Top pill
        pill_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf_pill = pill_box.text_frame
        tf_pill.word_wrap = True
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = f"●  LIFEQR REAL-WORLD ECOSYSTEM  |  {category.upper()}"
        p_pill.font.size = Pt(10)
        p_pill.font.bold = True
        p_pill.font.color.rgb = COLOR_RED

        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        if subtitle:
            p_sub = tf_title.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = COLOR_GRAY_MUTED

    # =========================================================================
    # SLIDE 1: Title Slide (Cover)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, COLOR_BG_DARK)

    # Decorative Border line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(0.15), Inches(4.8))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_RED
    line.line.fill.background()

    # Title box
    box1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.1), Inches(11.0), Inches(5.0))
    tf1 = box1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "LIFEQR CONNECTED HARDWARE ECOSYSTEM"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Real-World Emergency Triage, Hybrid NFC/QR Wearables & Dedicated Ambulance IoT Network"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_RED
    p2.space_before = Pt(12)

    p3 = tf1.add_paragraph()
    p3.text = "Solving the 'Golden Hour' medical data blackout through instant street-side hardware lookup,\nchild & rider accident protection, and automated ER trauma bay pre-notification."
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.space_before = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "PROJECT ROADMAP & HARDWARE SPECIFICATION  •  CONFIDENTIAL DECK"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GRAY_MUTED
    p4.space_before = Pt(50)

    # =========================================================================
    # SLIDE 2: The Critical Problem
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, COLOR_BG_DARK)
    add_header(slide2, "The Real-World Crisis", "The Street-Side Medical Data Blackout", "Why thousands of accident victims and students suffer preventable trauma deaths each year")

    cards_data = [
        ("70% Unconscious Victims", "Accident victims, seizing students, or stroke patients are physically unable to speak, unlock their phones, or communicate identity.", COLOR_RED),
        ("15-25 Min Golden Hour Lost", "Paramedics arrive blind without blood group, allergy warnings (penicillin), or chronic diseases (diabetes, hemophilia).", COLOR_AMBER),
        ("Paramedic Phone Reliance", "Responders currently rely on personal smartphones with slow autofocus, low battery, bad rain lighting, or broken glass.", COLOR_RED),
        ("Disconnected ER Desks", "Hospital trauma bays only know about critical emergencies when ambulance sirens physically pull into the hospital driveway.", COLOR_AMBER)
    ]

    for i, (head, desc, acc_col) in enumerate(cards_data):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.2)
        card = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.4))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_DARK
        card.line.color.rgb = acc_col
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = f"CRITICAL GAP 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = acc_col

        p2 = tf.add_paragraph()
        p2.text = head
        p2.font.size = Pt(17)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(10)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 3: The 3 Core Components Solution Architecture
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, COLOR_BG_DARK)
    add_header(slide3, "System Architecture", "The 3 Pillars of Connected LifeQR", "A unified, hardware-backed emergency response loop from citizen to trauma center")

    pillars = [
        ("1. CITIZEN & PATIENT", "Hybrid NFC/QR Wearables", "Zero-battery smart medical cards, rider helmet crash decals, and student backpack safety tags. Passive, waterproof, permanent.", COLOR_GREEN),
        ("2. PARAMEDIC & AMBULANCE", "Tactical Triage Terminal", "Ruggedized field hardware with instant dual NFC tap & 2D laser barcode scanner, GPS proximity detector, and 1-touch ER broadcaster.", COLOR_RED),
        ("3. DOCTOR & CLINICAL ER", "Desktop Hub & Live Stream", "Doctor USB/Bluetooth NFC scanner cradle + Hospital ER Triage Command Center reserving trauma bays before ambulance arrival.", COLOR_AMBER)
    ]

    for i, (badge, title, desc, col) in enumerate(pillars):
        x = Inches(0.8 + (i * 3.95))
        y = Inches(2.2)
        card = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.75), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_DARK
        card.line.color.rgb = col
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.35)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = badge
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(19)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 4: Component 1 — Patient NFC Card & Wearables (With Image)
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, COLOR_BG_DARK)
    add_header(slide4, "Component 01", "Citizen Wearables: Hybrid NFC & QR Physical IDs", "Zero-battery, tamper-resistant medical credentials for daily protection")

    # Left: Text Card
    left_card = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8))
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = COLOR_CARD_DARK
    left_card.line.color.rgb = COLOR_BORDER

    tf4 = left_card.text_frame
    tf4.word_wrap = True
    tf4.margin_left = Inches(0.35)
    tf4.margin_top = Inches(0.35)
    tf4.margin_right = Inches(0.35)

    bullets4 = [
        ("NFC Chip Embedded in PVC Card", "Passive 13.56 MHz NTAG213/424 chip read by any smartphone or responder hardware in under 0.2 seconds."),
        ("High-Contrast Laser QR", "Printed on the card surface for instant optical scanning when NFC is obstructed."),
        ("Biker & Helmet Crash Decals", "Retroreflective, weatherproof decals placed on rider helmets for immediate bystander or EMS access."),
        ("Student & Pediatric Bag Tags", "Clip-on tags for school bags protecting children with epilepsy, asthma, diabetes, or severe nut allergies."),
        ("Zero-Power & Lifetime Durability", "No charging required; waterproof, heat-resistant, and functional for 10+ years.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets4):
        p_t = tf4.add_paragraph() if i > 0 else tf4.paragraphs[0]
        p_t.text = f"✔  {b_title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_RED
        if i > 0: p_t.space_before = Pt(10)

        p_d = tf4.add_paragraph()
        p_d.text = b_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT

    # Right: Rendered Mockup Image
    img_path1 = os.path.join(assets_dir, 'nfc_card_wearables.jpg')
    if os.path.exists(img_path1):
        slide4.shapes.add_picture(img_path1, Inches(6.8), Inches(2.0), width=Inches(5.7), height=Inches(4.8))

    # =========================================================================
    # SLIDE 5: Component 2 — Ambulance Handheld Hardware (With Image)
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, COLOR_BG_DARK)
    add_header(slide5, "Component 02", "Ambulance Tactical Triage Terminal", "Dedicated field-grade hardware unit for paramedics & emergency crews")

    # Left: Mockup Image
    img_path2 = os.path.join(assets_dir, 'ambulance_triage_device.jpg')
    if os.path.exists(img_path2):
        slide5.shapes.add_picture(img_path2, Inches(0.8), Inches(2.0), width=Inches(5.7), height=Inches(4.8))

    # Right: Text Card
    right_card = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.7), Inches(2.0), Inches(5.8), Inches(4.8))
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = COLOR_CARD_DARK
    right_card.line.color.rgb = COLOR_RED
    right_card.line.width = Pt(1.5)

    tf5 = right_card.text_frame
    tf5.word_wrap = True
    tf5.margin_left = Inches(0.35)
    tf5.margin_top = Inches(0.35)
    tf5.margin_right = Inches(0.35)

    bullets5 = [
        ("Dual Sensor Engine", "Top NFC halo antenna + 2D CMOS barcode scanner reads cards, tags, and helmet stickers through grime or rain in < 0.2s."),
        ("Immediate Triage HUD", "Instantly projects Blood Group, Lethal Allergies, Current Medications, and Emergency Contacts on high-contrast screen."),
        ("Automated Hospital Detection", "Integrated GPS module detects nearest Level-1/2 Trauma Centers and calculates traffic-aware arrival countdown."),
        ("1-Touch ER Dispatch Stream", "Transmits live vitals (HR, SpO2, BP) and triage priority (CRITICAL/URGENT) directly to receiving hospital ER desk."),
        ("Rugged Industrial Design", "IP67 dust/waterproof, drop-resistant rubber bumper casing for harsh accident environments.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets5):
        p_t = tf5.add_paragraph() if i > 0 else tf5.paragraphs[0]
        p_t.text = f"★  {b_title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_RED
        if i > 0: p_t.space_before = Pt(10)

        p_d = tf5.add_paragraph()
        p_d.text = b_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT

    # =========================================================================
    # SLIDE 6: Component 3 — Doctor Desktop Hardware & ER Command (With Image)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, COLOR_BG_DARK)
    add_header(slide6, "Component 03", "Doctor Desktop Cradle & ER Command Center", "Zero-friction clinic intake and in-transit ambulance telemetry reception")

    # Left: Text Card
    left_card6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.8), Inches(4.8))
    left_card6.fill.solid()
    left_card6.fill.fore_color.rgb = COLOR_CARD_DARK
    left_card6.line.color.rgb = COLOR_BORDER

    tf6 = left_card6.text_frame
    tf6.word_wrap = True
    tf6.margin_left = Inches(0.35)
    tf6.margin_top = Inches(0.35)
    tf6.margin_right = Inches(0.35)

    bullets6 = [
        ("Plug-and-Play USB/BLE Cradle", "Medical-grade desktop reader sits beside doctor monitor. Doctor simply taps patient card or scans badge."),
        ("Zero-Typing Patient Intake", "Auto-populates patient profile, past medical surgeries, and active prescriptions without manual data entry errors."),
        ("Live In-Transit ER Radar", "Receiving ER staff monitors approaching ambulances, live ETA countdown, and paramedic vitals stream."),
        ("Trauma Bay Pre-Allocation", "Allows ER charge nurses to reserve Trauma Bay 1/2/3 and stage matched blood units before patient arrival."),
        ("AI Clinical Decision Support", "Assists physician with AI drug interaction alerts, differential diagnosis, and instant SOAP notes.")
    ]

    for i, (b_title, b_desc) in enumerate(bullets6):
        p_t = tf6.add_paragraph() if i > 0 else tf6.paragraphs[0]
        p_t.text = f"✔  {b_title}"
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_GREEN
        if i > 0: p_t.space_before = Pt(10)

        p_d = tf6.add_paragraph()
        p_d.text = b_desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT

    # Right: Mockup Image
    img_path3 = os.path.join(assets_dir, 'doctor_desktop_scanner.jpg')
    if os.path.exists(img_path3):
        slide6.shapes.add_picture(img_path3, Inches(6.8), Inches(2.0), width=Inches(5.7), height=Inches(4.8))

    # =========================================================================
    # SLIDE 7: Real-Life Incident Scenario (Timeline)
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7, COLOR_BG_DARK)
    add_header(slide7, "Real-World Workflow", "End-to-End Emergency Incident Timeline", "How a motorcycle highway crash is handled from impact to resuscitation")

    steps = [
        ("00:00", "Crash Occurs", "Motorcyclist unconscious on highway. Bystander calls emergency hotline."),
        ("00:04", "EMS Arrival", "Paramedic arrives with LifeQR tactical terminal in hand."),
        ("00:05", "1-Tap Scan", "Paramedic taps helmet QR decal / wallet NFC card. Patient ID resolved in 0.2s."),
        ("00:06", "Triage Displayed", "Terminal shows: Blood O+, Penicillin Anaphylaxis, Severe Asthma. Family ICE notified."),
        ("00:08", "Auto-ER Alert", "Device detects Nearest Hospital (City Trauma Care). Broadcasts vitals & 7m ETA."),
        ("00:15", "Handover Complete", "Patient arrives; Trauma Bay 1 pre-allocated, O-negative blood thawed. Life saved.")
    ]

    for i, (t_time, t_title, t_desc) in enumerate(steps):
        x = Inches(0.8 + (i * 1.95))
        y = Inches(2.3)
        node = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(1.8), Inches(4.3))
        node.fill.solid()
        node.fill.fore_color.rgb = COLOR_CARD_DARK
        node.line.color.rgb = COLOR_RED if i in [2, 4, 5] else COLOR_BORDER
        node.line.width = Pt(1.5)

        tf = node.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = t_time
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_RED

        p2 = tf.add_paragraph()
        p2.text = t_title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(8)

        p3 = tf.add_paragraph()
        p3.text = t_desc
        p3.font.size = Pt(10)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(10)

    # =========================================================================
    # SLIDE 8: High-Risk Demographics & Safety Use Cases
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8, COLOR_BG_DARK)
    add_header(slide8, "Target Demographics", "High-Impact Real-Life Use Cases", "Deploying LifeQR cards and decals across vulnerable population groups")

    use_cases = [
        ("Two-Wheeler & Highway Bikers", "High crash risk. Helmet QR decal and wallet card provide instant triage when rider is unconscious with protective gear on.", "Rider Helmet Decal & Wallet Card", COLOR_RED),
        ("School & College Students", "Protects children with chronic conditions (epilepsy, asthma, severe allergies). If student collapses, teachers scan backpack tag immediately.", "Backpack NFC Badge & ID Clip", COLOR_AMBER),
        ("Senior Citizens & Dementia Care", "Elderly individuals prone to falls, memory loss, or wandering. Bracelet or card alerts bystanders and contacts adult children instantly.", "Waterproof Medical Wristband", COLOR_GREEN),
        ("Diabetic & Hemophilic Patients", "Critical medical condition tags preventing lethal administration of standard emergency drugs by first responders.", "Laser-Engraved Metal NFC Card", COLOR_RED)
    ]

    for i, (title, desc, form_factor, col) in enumerate(use_cases):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.2 + (row * 2.4))
        
        card = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD_DARK
        card.line.color.rgb = col
        card.line.width = Pt(1.5)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = f"Form Factor: {form_factor}"
        p3.font.size = Pt(10)
        p3.font.bold = True
        p3.font.color.rgb = col
        p3.space_before = Pt(6)

    # =========================================================================
    # SLIDE 9: Manufacturing & Distribution Business Model
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9, COLOR_BG_DARK)
    add_header(slide9, "Business & Distribution", "Sustainable Production & Distribution Strategy", "Low manufacturing cost, high societal impact, and self-sustaining revenue")

    model_cols = [
        ("NFC Card Production", "Low Unit Cost", "• High-volume PVC/PET NFC Card: ~$0.30 - $0.50 per unit.\n• Weatherproof 3M Reflective Decal: ~$0.15 per unit.\n• Sold at nominal price ($2 - $3) at partner clinics, pharmacies, and driving license centers.", COLOR_GREEN),
        ("Hospital & School Supply", "B2B Institutional Supply", "• Supplied in bulk to school districts for student IDs.\n• Bundled with two-wheeler vehicle delivery at motorcycle dealerships.\n• Distributed at hospital discharge for chronic patients.", COLOR_AMBER),
        ("Hardware-as-a-Service", "Ambulance Fleet Model", "• Handheld triage terminals leased to private and municipal ambulance services ($25/mo per unit).\n• Free doctor desktop software & affordable USB scanner cradles ($35 one-time).", COLOR_RED),
        ("Value-Added Services", "Emergency Platform", "• Automated SMS/WhatsApp multi-contact dispatch.\n• Secure cloud profile syncing and AI prescription safety check.\n• Corporate and campus safety subscription packages.", COLOR_GREEN)
    ]

    for i, (head, subhead, body, col) in enumerate(model_cols):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.2)
        c = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.5))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = head
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = subhead
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = col
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 10: Technical Hardware Specifications
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10, COLOR_BG_DARK)
    add_header(slide10, "Engineering Specs", "Hardware Architecture & IoT Specifications", "Component-level breakdown for manufacturing the LifeQR terminal and wearables")

    spec_blocks = [
        ("Ambulance Triage Terminal", [
            ("Core SoC", "ESP32-S3 Dual-Core Xtensa LX7 (240MHz) with hardware cryptographic engine"),
            ("NFC Reader", "NXP PN532 transceiver supporting ISO14443A/B, MIFARE, and FeliCa (13.56MHz)"),
            ("Barcode Scanner", "2D CMOS High-Speed Optical Imager with red aiming crosshair and white fill LED"),
            ("Connectivity", "Quectel EC200U 4G LTE-M module with integrated GNSS GPS/GLONASS positioning"),
            ("Display & Battery", "3.5-inch High-Brightness Sunlight-Readable LCD; 4000mAh Li-ion (18h field duty)")
        ]),
        ("Doctor Desktop Scanner Cradle", [
            ("Interface", "High-speed USB 2.0 Type-C + Bluetooth 5.2 Low Energy"),
            ("Sensor", "Contactless NFC flat-bed reader + 120fps overhead QR camera sensor"),
            ("OS Compatibility", "Windows 10/11, macOS, Linux, Android Hospital Workstation plug-and-play"),
            ("Housing", "Antimicrobial ABS medical-grade plastic chassis resistant to hospital disinfectants")
        ])
    ]

    for i, (title, specs) in enumerate(spec_blocks):
        x = Inches(0.8 + (i * 5.9))
        y = Inches(2.2)
        b = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(4.5))
        b.fill.solid()
        b.fill.fore_color.rgb = COLOR_CARD_DARK
        b.line.color.rgb = COLOR_BORDER

        tf = b.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_RED

        for s_label, s_val in specs:
            p_l = tf.add_paragraph()
            p_l.text = f"•  {s_label}:"
            p_l.font.size = Pt(12)
            p_l.font.bold = True
            p_l.font.color.rgb = COLOR_WHITE
            p_l.space_before = Pt(8)

            p_v = tf.add_paragraph()
            p_v.text = f"   {s_val}"
            p_v.font.size = Pt(11)
            p_v.font.color.rgb = COLOR_GRAY_LIGHT

    # =========================================================================
    # SLIDE 11: Real-World Implementation Roadmap
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11, COLOR_BG_DARK)
    add_header(slide11, "Execution Roadmap", "Phased Commercial & Field Deployment", "From hardware prototyping to municipal emergency system adoption")

    phases = [
        ("Phase 1: Prototyping", "Months 1-3", "• Finalize ESP32 + PN532 + 2D scanner PCB design.\n• 3D print rugged enclosure.\n• Produce initial batch of 500 NFC smart cards & helmet decals.\n• Validate firmware with LifeQR cloud API.", COLOR_AMBER),
        ("Phase 2: EMS Pilot", "Months 4-6", "• Field pilot with 20 municipal ambulance units.\n• Deploy 5 doctor desktop cradles in Level-1 Trauma ER.\n• Real-world crash response testing & paramedic UX tuning.\n• Establish average triage latency benchmark (< 5 sec).", COLOR_RED),
        ("Phase 3: Community Rollout", "Months 7-10", "• Partner with motorcycle riding clubs for helmet decal distribution.\n• School district safety tag program for 10,000 students.\n• Partner with private hospital networks for patient onboarding.", COLOR_GREEN),
        ("Phase 4: Scale & Gov Integration", "Months 11+", "• Integrate with national emergency dispatch hotlines (911 / 112 / 108).\n• Commercial mass manufacturing of triage hardware.\n• Expansion into statewide highway emergency networks.", COLOR_WHITE)
    ]

    for i, (phase_title, phase_time, phase_body, col) in enumerate(phases):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.2)
        c = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.5))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.3)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = phase_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = phase_time
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = col
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = phase_body
        p3.font.size = Pt(11)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 12: Conclusion & Call to Action
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12, COLOR_BG_DARK)

    # Central Impact Card
    cta_card = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.2), Inches(10.33), Inches(5.1))
    cta_card.fill.solid()
    cta_card.fill.fore_color.rgb = COLOR_CARD_DARK
    cta_card.line.color.rgb = COLOR_RED
    cta_card.line.width = Pt(2)

    tf12 = cta_card.text_frame
    tf12.word_wrap = True
    tf12.margin_left = Inches(0.6)
    tf12.margin_top = Inches(0.6)
    tf12.margin_right = Inches(0.6)

    p = tf12.paragraphs[0]
    p.text = "WHEN SECONDS COUNT, LIFEQR SAVES LIVES."
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf12.add_paragraph()
    p2.text = "Bridging Physical Reality and Emergency Healthcare"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_RED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)

    p3 = tf12.add_paragraph()
    p3.text = "By uniting passive hybrid NFC/QR wearables, dedicated ambulance triage terminals,\nand real-time hospital notification networks, LifeQR eliminates the deadly 'Golden Hour' data void.\nWe are ready to build the physical devices, distribute the wearables, and transform accident survival."
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    p4 = tf12.add_paragraph()
    p4.text = "CONTACT & PILOT INQUIRIES: info@lifeqr.org  •  LIFEQR HARDWARE LABS"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GREEN
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(35)

    # Save presentation
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/LifeQR_RealWorld_Ecosystem.pptx'))
    prs.save(output_path)
    print(f"[Success] Presentation generated at: {output_path}")

if __name__ == '__main__':
    create_presentation()
