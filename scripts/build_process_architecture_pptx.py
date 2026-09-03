"""
LifeQR Process Architecture Master PowerPoint Generator
Generates a 16-slide, 16:9 widescreen presentation deck using python-pptx.
Incorporates high-resolution dedicated architecture flowcharts for:
1. Patient & Citizen Lifecycle Architecture (Registration, NTAG 424 DNA Encoding, Laser QR, Wearables, Cloud Vault)
2. Highway Motorcycle Accident Triage Architecture (Crash, IMU Sensor, Tactical Terminal, Dual Scan, Offline EEPROM, GPS Dispatch)
3. Hospital & Clinical ER Resuscitation Architecture (In-Transit Telemetry, Trauma Bay 1 Reservation, Blood Bank Thawing, Doctor Desktop Cradle, FHIR v4 EHR Integration)
4. Full-Stack IoT Ecosystem Architecture & Tactical Terminal Hardware Exploded Blueprint.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_process_architecture_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors (Deep Cybernetic Slate, Signal Red & Neon Cyan)
    COLOR_BG_DARK = RGBColor(11, 14, 19)       # #0b0e13
    COLOR_CARD_DARK = RGBColor(20, 24, 33)     # #141821
    COLOR_BORDER = RGBColor(45, 52, 66)        # #2d3442
    COLOR_WHITE = RGBColor(255, 255, 255)      # #ffffff
    COLOR_RED = RGBColor(225, 29, 46)          # #E11D2E
    COLOR_CYAN = RGBColor(14, 165, 233)        # #0ea5e9
    COLOR_GREEN = RGBColor(16, 185, 129)       # #10b981
    COLOR_AMBER = RGBColor(245, 158, 11)       # #f59e0b
    COLOR_GRAY_LIGHT = RGBColor(209, 213, 219) # #d1d5db
    COLOR_GRAY_MUTED = RGBColor(156, 163, 175) # #9ca3af

    assets_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/assets'))

    def set_slide_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=None):
        pill_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_pill = pill_box.text_frame
        tf_pill.word_wrap = True
        p_pill = tf_pill.paragraphs[0]
        p_pill.text = f"●  LIFEQR PROCESS ARCHITECTURE MASTER  |  {category.upper()}"
        p_pill.font.size = Pt(10)
        p_pill.font.bold = True
        p_pill.font.color.rgb = COLOR_RED

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.7), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(25)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE

        if subtitle:
            p_sub = tf_title.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(12)
            p_sub.font.color.rgb = COLOR_GRAY_MUTED

    # =========================================================================
    # SLIDE 1: Title Slide (Cover)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1, COLOR_BG_DARK)

    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(0.18), Inches(5.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_RED
    bar.line.fill.background()

    box1 = slide1.shapes.add_textbox(Inches(1.25), Inches(1.1), Inches(11.0), Inches(5.2))
    tf1 = box1.text_frame
    tf1.word_wrap = True

    p = tf1.paragraphs[0]
    p.text = "LIFEQR PROCESS ARCHITECTURE MASTER"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p2 = tf1.add_paragraph()
    p2.text = "Comprehensive Process Blueprints for Patient Lifecycle, Accident Triage & Hospital Resuscitation"
    p2.font.size = Pt(17)
    p2.font.color.rgb = COLOR_RED
    p2.space_before = Pt(10)

    p3 = tf1.add_paragraph()
    p3.text = "Step-by-step technical architectures with dedicated engineering flowcharts:\n• Process 1: Patient Registration, NTAG 424 Cryptographic Encoding & Wearable Deployment\n• Process 2: Highway Motorcycle Crash, IMU Shock Sensing & Tactical Paramedic Triage\n• Process 3: Real-Time In-Transit Hospital Telemetry, Trauma Bay Allocation & Doctor Desktop Hub"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.space_before = Pt(20)

    p4 = tf1.add_paragraph()
    p4.text = "PROCESS SPECIFICATIONS  •  VISUAL ARCHITECTURE BLUEPRINTS  •  REAL-WORLD WORKFLOWS"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GRAY_MUTED
    p4.space_before = Pt(45)

    # =========================================================================
    # SLIDE 2: Master Ecosystem Process Topology (With Image)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2, COLOR_BG_DARK)
    add_header(slide2, "Master Topology", "Full-Stack Connected Emergency Ecosystem", "The complete synchronized data loop connecting patient, scene, ambulance, and trauma bay")

    img_top = os.path.join(assets_dir, 'lifeqr_system_architecture.jpg')
    if os.path.exists(img_top):
        slide2.shapes.add_picture(img_top, Inches(0.8), Inches(1.9), width=Inches(7.2), height=Inches(5.0))

    card2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), Inches(1.9), Inches(4.3), Inches(5.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = COLOR_CARD_DARK
    card2.line.color.rgb = COLOR_BORDER

    tf2 = card2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.25)
    tf2.margin_top = Inches(0.25)
    tf2.margin_right = Inches(0.25)

    p = tf2.paragraphs[0]
    p.text = "THE 3 CORE PROCESS TIERS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    tiers_desc = [
        ("Process Tier 1: Patient Lifecycle", "Digital registration, NTAG 424 DNA programming, laser-etched QR card, helmet decal, and school student tag.", COLOR_GREEN),
        ("Process Tier 2: Accident & EMS Triage", "Crash event, IMU shock trigger, paramedic tactical dual-scan, offline EEPROM vitals read, and GPS hospital dispatch.", COLOR_RED),
        ("Process Tier 3: Hospital & Doctor ER", "Live telemetry reception, Trauma Bay 1 allocation, blood bank staging, doctor USB cradle intake, and FHIR EHR sync.", COLOR_CYAN)
    ]

    for t_head, t_body, t_col in tiers_desc:
        p_h = tf2.add_paragraph()
        p_h.text = f"■  {t_head}"
        p_h.font.size = Pt(12)
        p_h.font.bold = True
        p_h.font.color.rgb = t_col
        p_h.space_before = Pt(12)

        p_b = tf2.add_paragraph()
        p_b.text = t_body
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = COLOR_GRAY_LIGHT
        p_b.space_before = Pt(4)

    # =========================================================================
    # SLIDE 3: Process 1 Blueprint — Patient Lifecycle (With Image)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3, COLOR_BG_DARK)
    add_header(slide3, "Process Architecture 01", "Patient & Citizen Lifecycle Blueprint", "Registration, NTAG 424 DNA chip encoding, laser QR printing, and zero-knowledge encryption")

    img_pat = os.path.join(assets_dir, 'patient_process_arch.jpg')
    if os.path.exists(img_pat):
        slide3.shapes.add_picture(img_pat, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    card3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    card3.fill.solid()
    card3.fill.fore_color.rgb = COLOR_CARD_DARK
    card3.line.color.rgb = COLOR_GREEN
    card3.line.width = Pt(1.5)

    tf3 = card3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = Inches(0.2)
    tf3.margin_top = Inches(0.25)
    tf3.margin_right = Inches(0.2)

    p = tf3.paragraphs[0]
    p.text = "PROCESS 1 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_GREEN

    p_steps = [
        ("Step 1: Emergency Intake", "Patient inputs blood group, allergies, chronic ailments, and primary ICE phone contacts."),
        ("Step 2: NTAG 424 DNA", "Card issuance encoder writes rolling AES-128 SUN keys to prevent counterfeiting."),
        ("Step 3: Multi-Wearables", "Wallet smart card, rider helmet reflective decal, and student backpack tag issued."),
        ("Step 4: Encrypted Vault", "Data secured via AES-256 GCM client-side encryption with zero-knowledge keys.")
    ]

    for s_title, s_desc in p_steps:
        p_h = tf3.add_paragraph()
        p_h.text = f"✔ {s_title}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(10)

        p_d = tf3.add_paragraph()
        p_d.text = s_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(3)

    # =========================================================================
    # SLIDE 4: Patient Lifecycle Deep Dive: Step-by-Step Technical Flow
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4, COLOR_BG_DARK)
    add_header(slide4, "Patient Process Deep-Dive", "Technical Lifecycle: Registration to Card Issuance", "Under the hood of the patient medical onboarding and cryptographic encoding pipeline")

    patient_phases = [
        ("1. Digital Medical Onboarding", "• Verified Identity & Consent\n• Critical Vitals Intake: Blood Group (A+, B+, O+, AB-), Severe Allergies (Penicillin, Nuts), Chronic Conditions (Diabetes, Epilepsy).\n• ICE Contacts: Primary, secondary, family doctor.\n• Emergency Directives & DNR flags.", COLOR_GREEN),
        ("2. NTAG 424 DNA Encoding", "• Issuance terminal programs NXP NTAG 424 DNA.\n• Generates Master Application Key (K0-K4).\n• Configures Secure Unique NFC (SUN) with AES-128 CMAC.\n• Programs 512-byte encrypted offline EEPROM payload for zero-internet emergency access.", COLOR_CYAN),
        ("3. Laser Etching & Surface Printing", "• Micro-QR code laser-engraved onto composite PVC/PET surface.\n• Encoded with Level H error correction (withstands 30% physical abrasion or blood).\n• Distinct high-contrast Swiss Brutalist emergency layout with patient unique alphanumeric ID.", COLOR_AMBER),
        ("4. Zero-Knowledge Cloud Storage", "• Medical records encrypted client-side using AES-256 GCM.\n• Master keys stored in hardware security module (HSM).\n• Role-based dynamic data masking (Public bystander vs paramedic vs physician).\n• Immutable access log generated on each tap.", COLOR_RED)
    ]

    for i, (p_title, p_body, p_col) in enumerate(patient_phases):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = p_col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"PHASE 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = p_col

        p2 = tf.add_paragraph()
        p2.text = p_title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = p_body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 5: School & Pediatric Chronic Care Architecture
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5, COLOR_BG_DARK)
    add_header(slide5, "Community Protection", "School District & Pediatric Chronic Safety Architecture", "Protecting students with epilepsy, severe asthma, type-1 diabetes, and nut anaphylaxis")

    school_steps = [
        ("Step 1: School District Portal", "Campus registered nurses and administrators enroll students with documented chronic conditions via the LifeQR School Safety Dashboard.", COLOR_CYAN),
        ("Step 2: Bag Tag & Wristband Issuance", "Children receive durable silicone wristbands and clip-on backpack tags equipped with dual NFC and laser QR credentials.", COLOR_GREEN),
        ("Step 3: Sudden Classroom / Field Incident", "If a student faints or experiences an epileptic seizure on the playground, teachers or bus drivers tap the badge in 1.8s.", COLOR_AMBER),
        ("Step 4: Instant Protocol & Parent Alert", "Displays immediate emergency instructions ('Administer EpiPen from left pouch; do not restrain') and dispatches automated SMS/call to parents.", COLOR_RED)
    ]

    for i, (s_title, s_desc, s_col) in enumerate(school_steps):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = s_col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = s_title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = s_desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 6: Process 2 Blueprint — Highway Motorcycle Accident Triage (With Image)
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6, COLOR_BG_DARK)
    add_header(slide6, "Process Architecture 02", "Highway Motorcycle Accident Triage Blueprint", "IMU crash detection, paramedic tactical dual-scan, offline EEPROM decode & GPS dispatch")

    img_acc = os.path.join(assets_dir, 'accident_process_arch.jpg')
    if os.path.exists(img_acc):
        slide6.shapes.add_picture(img_acc, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    card6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    card6.fill.solid()
    card6.fill.fore_color.rgb = COLOR_CARD_DARK
    card6.line.color.rgb = COLOR_RED
    card6.line.width = Pt(1.5)

    tf6 = card6.text_frame
    tf6.word_wrap = True
    tf6.margin_left = Inches(0.2)
    tf6.margin_top = Inches(0.25)
    tf6.margin_right = Inches(0.2)

    p = tf6.paragraphs[0]
    p.text = "PROCESS 2 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    acc_steps = [
        ("1. Motorcycle Crash", "IoT IMU sensor beacon detects >15G impact, broadcasting automated BLE emergency beacon."),
        ("2. Paramedic Arrives", "Arrives with rugged tactical handheld terminal; no personal smartphone delays."),
        ("3. Instant 0.2s Scan", "Scans helmet QR decal or taps wallet NFC card; decrypts O+ blood and allergy offline."),
        ("4. GPS CAD Dispatch", "Auto-identifies nearest Level-1 trauma center & transmits live vitals en route.")
    ]

    for a_title, a_desc in acc_steps:
        p_h = tf6.add_paragraph()
        p_h.text = f"★ {a_title}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(10)

        p_d = tf6.add_paragraph()
        p_d.text = a_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(3)

    # =========================================================================
    # SLIDE 7: Motorcycle Crash & Autonomous Detection Flow
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7, COLOR_BG_DARK)
    add_header(slide7, "Crash Detection Mechanics", "Autonomous Motorcycle Crash & E-Call Sensor Flow", "How the LifeQR IoT helmet beacon identifies high-G crashes and eliminates delay")

    c_steps = [
        ("1. Impact Event (>15G)", "Rider strikes barrier or vehicle. 6-axis Bosch BMI270 accelerometer detects severe deceleration spike (>15G) combined with high angular velocity.", COLOR_RED),
        ("2. 30s Stillness Verification", "Algorithm verifies post-crash victim stillness for 30 seconds to filter out benign drops (e.g. helmet falling from seat), preventing false 911 calls.", COLOR_AMBER),
        ("3. Autonomous BLE/E-Call Broadcast", "Beacon triggers high-power BLE 5.3 emergency advertisement packet containing the rider's LifeQR ID and GPS coordinates to passing phones & CAD.", COLOR_CYAN),
        ("4. Bystander / First Responder Scan", "When responders or bystanders reach the victim, the 3M retroreflective helmet decal is scanned from 1 meter away without removing the helmet.", COLOR_GREEN)
    ]

    for i, (head, body, col) in enumerate(c_steps):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"STEP 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = col

        p2 = tf.add_paragraph()
        p2.text = head
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 8: Ambulance Tactical Terminal Hardware Workflow (With Image)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8, COLOR_BG_DARK)
    add_header(slide8, "Ambulance Hardware Workflow", "Paramedic Tactical Terminal: Field Execution Flow", "Zero-latency dual scanning, offline EEPROM decoding, and live vitals telemetry")

    img_dev = os.path.join(assets_dir, 'ambulance_triage_device.jpg')
    if os.path.exists(img_dev):
        slide8.shapes.add_picture(img_dev, Inches(0.8), Inches(1.9), width=Inches(6.6), height=Inches(4.9))

    card8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.9), Inches(4.9))
    card8.fill.solid()
    card8.fill.fore_color.rgb = COLOR_CARD_DARK
    card8.line.color.rgb = COLOR_RED
    card8.line.width = Pt(1.5)

    tf8 = card8.text_frame
    tf8.word_wrap = True
    tf8.margin_left = Inches(0.25)
    tf8.margin_top = Inches(0.25)
    tf8.margin_right = Inches(0.25)

    p = tf8.paragraphs[0]
    p.text = "FIELD RESCUE WORKFLOW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED

    dev_points = [
        ("Sub-Second Dual Sensor Scan", "Paramedic approaches unconscious victim; taps card on top NFC halo or scans helmet QR decal. Sensor responds in < 200 milliseconds even in heavy rain."),
        ("Offline EEPROM Vitals Decrypt", "If cellular reception is 0%, the terminal decrypts the 512-byte payload stored in the card's chip. Projects Blood Type O+ and Penicillin Allergy on screen."),
        ("1-Touch Vitals Entry & Stream", "Paramedic attaches pulse oximeter and cuff; taps HR (134), SpO2 (89%), BP (95/58). Data streams continuously to hospital ER over 4G/Satellite."),
        ("High-Contrast Triage HUD", "Sunlight-readable 800-nit screen displays Red (Urgent) triage banner, preventing lethal contraindications in noisy, chaotic crash sites.")
    ]

    for d_head, d_body in dev_points:
        p_h = tf8.add_paragraph()
        p_h.text = f"★ {d_head}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_b = tf8.add_paragraph()
        p_b.text = d_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = COLOR_GRAY_LIGHT
        p_b.space_before = Pt(3)

    # =========================================================================
    # SLIDE 9: In-Transit Hospital Proximity & Auto-Notification Flow
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide9, COLOR_BG_DARK)
    add_header(slide9, "Dispatch Flow", "Proximity GPS & Hospital Auto-Notification Architecture", "Automated trauma center selection, ETA calculation, and advance bay reservation")

    flow_nodes = [
        ("1. GPS Spatial Query", "Quectel GNSS module fixes exact crash coordinates and queries the regional trauma hospital registry within 25 km radius.", COLOR_CYAN),
        ("2. Capability Matching", "Evaluates live hospital telemetry: verified on-duty trauma surgeon, CT scanner availability, and O-negative blood stock.", COLOR_GREEN),
        ("3. Automated Reservation", "Transmits an authenticated reservation token to City Trauma Care ER. Locks Trauma Bay 1 and thaws blood units 7 minutes prior to arrival.", COLOR_AMBER),
        ("4. Live In-Transit Radar", "Ambulance streams continuous GPS coordinates, speed, traffic-adjusted ETA, and patient vitals to the receiving hospital command desk.", COLOR_RED)
    ]

    for i, (f_title, f_desc, f_col) in enumerate(flow_nodes):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = f_col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"PHASE 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = f_col

        p2 = tf.add_paragraph()
        p2.text = f_title
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = f_desc
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 10: Process 3 Blueprint — Hospital ER Resuscitation (With Image)
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide10, COLOR_BG_DARK)
    add_header(slide10, "Process Architecture 03", "Hospital & Clinical ER Resuscitation Blueprint", "In-transit telemetry, automated trauma bay reservation, doctor desktop cradle & FHIR sync")

    img_hosp = os.path.join(assets_dir, 'hospital_process_arch.jpg')
    if os.path.exists(img_hosp):
        slide10.shapes.add_picture(img_hosp, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    card10 = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    card10.fill.solid()
    card10.fill.fore_color.rgb = COLOR_CARD_DARK
    card10.line.color.rgb = COLOR_CYAN
    card10.line.width = Pt(1.5)

    tf10 = card10.text_frame
    tf10.word_wrap = True
    tf10.margin_left = Inches(0.2)
    tf10.margin_top = Inches(0.25)
    tf10.margin_right = Inches(0.2)

    p = tf10.paragraphs[0]
    p.text = "PROCESS 3 HIGHLIGHTS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    hosp_steps = [
        ("1. In-Transit Vitals", "ER monitor receives live countdown (ETA 07:12) and vitals (HR 134, SpO2 89%, BP 95/58)."),
        ("2. Bay 1 & Blood Staged", "Trauma Bay 1 locked; blood bank thaws matched O-negative units before siren arrives."),
        ("3. Doctor Desktop Cradle", "Doctor taps card on USB/BLE scanner cradle; profile populates with zero manual typing."),
        ("4. FHIR & AI Safety", "Auto-syncs into Epic/Cerner EHR; flags penicillin allergy and generates SOAP chart.")
    ]

    for h_title, h_desc in hosp_steps:
        p_h = tf10.add_paragraph()
        p_h.text = f"✔ {h_title}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(10)

        p_d = tf10.add_paragraph()
        p_d.text = h_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(3)

    # =========================================================================
    # SLIDE 11: ER Trauma Bay Pre-Allocation & Advance Resuscitation Staging
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide11, COLOR_BG_DARK)
    add_header(slide11, "Trauma Bay Preparation", "Advance Resuscitation & Blood Bank Protocol", "Converting chaotic hospital arrival into an orchestrated clinical resuscitation")

    er_phases = [
        ("Advance Notification (ETA -10m)", "ER charge nurse receives siren alert and telemetry stream on central command console. Victim identified as 28yo male, motorcycle polytrauma, GCS 9.", COLOR_AMBER),
        ("Trauma Bay Lockout (ETA -7m)", "Trauma Bay 1 is electronically reserved on the board. Nursing team pre-sets rapid infuser, ventilator circuit, and emergency chest tube tray.", COLOR_RED),
        ("Blood Bank Staging (ETA -5m)", "Hospital blood refrigerator receives automated reservation: thaws 2 units of uncrossed O-negative PRBCs and 2 units of plasma. Ready at bay bedside.", COLOR_CYAN),
        ("Immediate Handover (ETA 0m)", "Ambulance backs into trauma dock; victim moved directly into Bay 1. Zero check-in registration delay. Resuscitation begins immediately.", COLOR_GREEN)
    ]

    for i, (title, desc, col) in enumerate(er_phases):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = slide11.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 12: Doctor Desktop Scanner Cradle & Zero-Typing Intake (With Image)
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide12, COLOR_BG_DARK)
    add_header(slide12, "Clinical Intake Architecture", "Doctor Desktop Scanner & Zero-Typing Consultation Flow", "Eliminating transcription errors, medication mix-ups, and consultation paperwork delays")

    img_doc = os.path.join(assets_dir, 'doctor_desktop_scanner.jpg')
    if os.path.exists(img_doc):
        slide12.shapes.add_picture(img_doc, Inches(0.8), Inches(1.9), width=Inches(6.6), height=Inches(4.9))

    card12 = slide12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.6), Inches(1.9), Inches(4.9), Inches(4.9))
    card12.fill.solid()
    card12.fill.fore_color.rgb = COLOR_CARD_DARK
    card12.line.color.rgb = COLOR_CYAN
    card12.line.width = Pt(1.5)

    tf12 = card12.text_frame
    tf12.word_wrap = True
    tf12.margin_left = Inches(0.25)
    tf12.margin_top = Inches(0.25)
    tf12.margin_right = Inches(0.25)

    p = tf12.paragraphs[0]
    p.text = "CLINICAL DESK ARCHITECTURE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    doc_points = [
        ("Plug-and-Play USB HID/CCID", "Medical-grade scanner cradle connects via USB Type-C to doctor's PC. Functions driverless as a smart card reader and keyboard wedge."),
        ("1-Tap Patient Record Ingestion", "Doctor simply taps the patient's card or scans QR code $\rightarrow$ complete clinical profile, past surgeries, and active prescriptions populate with zero manual typing."),
        ("AI Clinical Safety Copilot", "Integrated AI cross-references doctor's newly prescribed medication against patient's allergies and existing drugs, instantly flagging lethal drug interactions."),
        ("Automated SOAP Note Generation", "Generates compliant Subjective, Objective, Assessment, Plan (SOAP) clinical encounter documentation ready for physician signature.")
    ]

    for d_head, d_body in doc_points:
        p_h = tf12.add_paragraph()
        p_h.text = f"✔ {d_head}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(8)

        p_b = tf12.add_paragraph()
        p_b.text = d_body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = COLOR_GRAY_LIGHT
        p_b.space_before = Pt(3)

    # =========================================================================
    # SLIDE 13: Clinical Interoperability Architecture (HL7 & FHIR v4)
    # =========================================================================
    slide13 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide13, COLOR_BG_DARK)
    add_header(slide13, "EHR Interoperability", "HL7 & FHIR v4 Clinical Integration Pipeline", "Standardized JSON health exchange connecting LifeQR with global hospital systems")

    fhir_blocks = [
        ("FHIR Patient Resource", "Maps legal name, MRN, date of birth, emergency contacts, and blood group according to US Core & international HL7 standards.", COLOR_CYAN),
        ("FHIR AllergyIntolerance", "Directly transmits critical anaphylactic risks (penicillin, sulfa, latex, nuts) with verification status: 'confirmed' and criticality: 'high'.", COLOR_RED),
        ("FHIR Condition & History", "Standardized SNOMED-CT and ICD-10 medical condition encoding (e.g. Type-1 Diabetes E10.9, Severe Asthma J45.9).", COLOR_AMBER),
        ("FHIR Encounter & Handover", "Encapsulates paramedic field vitals (HR, SpO2, BP), Glasgow Coma Scale, and incident timestamp into an active emergency encounter record.", COLOR_GREEN)
    ]

    for i, (f_head, f_body, f_col) in enumerate(fhir_blocks):
        x = Inches(0.8 + (i * 2.95))
        y = Inches(2.0)
        c = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.75), Inches(4.8))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = f_col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"RESOURCE 0{i+1}"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = f_col

        p2 = tf.add_paragraph()
        p2.text = f_head
        p2.font.size = Pt(15)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_WHITE
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = f_body
        p3.font.size = Pt(10.5)
        p3.font.color.rgb = COLOR_GRAY_LIGHT
        p3.space_before = Pt(12)

    # =========================================================================
    # SLIDE 14: Security, Privacy & Zero-Trust Governance
    # =========================================================================
    slide14 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide14, COLOR_BG_DARK)
    add_header(slide14, "Security & Compliance", "Zero-Trust Privacy & Multi-Tiered Access Governance", "Complying with HIPAA, GDPR, and medical confidentiality standards")

    security_points = [
        ("Tier A: Public / Bystander", "Scanning with a smartphone camera shows only life-critical data: Blood Group, Emergency ICE Contacts (1-tap dial), and severe anaphylaxis warnings. Zero personal addresses, social security numbers, or financial details are exposed.", COLOR_GREEN),
        ("Tier B: First Responder (EMS)", "Paramedic authenticated terminal unlocks detailed clinical profile: active medications, chronic ailments, cardiac history, and specialized emergency physician directives.", COLOR_AMBER),
        ("Tier C: Verified Hospital Physician", "Doctor desktop terminal unlocks full clinical consultation notes, past diagnostic lab reports, AI differential insights, and prescription safety validator history.", COLOR_CYAN),
        ("Cryptographic Audit Trail", "Every tap or scan is cryptographically signed, timestamped, and geotagged. Immutable audit logs recorded in compliance with HIPAA Section 164.312.", COLOR_RED)
    ]

    for i, (title, desc, col) in enumerate(security_points):
        row = i // 2
        col_idx = i % 2
        x = Inches(0.8 + (col_idx * 5.9))
        y = Inches(2.1 + (row * 2.4))

        c = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.6), Inches(2.1))
        c.fill.solid()
        c.fill.fore_color.rgb = COLOR_CARD_DARK
        c.line.color.rgb = col
        c.line.width = Pt(1.5)

        tf = c.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.25)
        tf.margin_right = Inches(0.25)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_GRAY_LIGHT
        p2.space_before = Pt(8)

    # =========================================================================
    # SLIDE 15: Hardware Engineering Blueprint (With Image)
    # =========================================================================
    slide15 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide15, COLOR_BG_DARK)
    add_header(slide15, "Hardware Schematics", "Tactical Triage Terminal: Component Architecture", "Exploded assembly breakdown for field-grade manufacturing and IP68 certification")

    img_exp2 = os.path.join(assets_dir, 'hardware_exploded_view.jpg')
    if os.path.exists(img_exp2):
        slide15.shapes.add_picture(img_exp2, Inches(0.8), Inches(1.85), width=Inches(8.2), height=Inches(5.1))

    card15 = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.2), Inches(1.85), Inches(3.3), Inches(5.1))
    card15.fill.solid()
    card15.fill.fore_color.rgb = COLOR_CARD_DARK
    card15.line.color.rgb = COLOR_AMBER
    card15.line.width = Pt(1.5)

    tf15 = card15.text_frame
    tf15.word_wrap = True
    tf15.margin_left = Inches(0.2)
    tf15.margin_top = Inches(0.25)
    tf15.margin_right = Inches(0.2)

    p = tf15.paragraphs[0]
    p.text = "HARDWARE SUBSYSTEMS"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_AMBER

    subsystems_desc = [
        ("Core SoC", "ESP32-S3 Dual-Core Xtensa @ 240MHz with hardware crypto engine."),
        ("Dual Sensors", "NXP PN532 NFC reader + 60fps 2D CMOS optical barcode imager."),
        ("Comms & GPS", "Quectel EC200U 4G LTE-M with GNSS GPS/GLONASS positioning."),
        ("Battery System", "4000mAh Li-Po (18h continuous duty) + USB-C PD & Qi wireless dock.")
    ]

    for s_name, s_det in subsystems_desc:
        p_h = tf15.add_paragraph()
        p_h.text = f"• {s_name}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = COLOR_WHITE
        p_h.space_before = Pt(10)

        p_d = tf15.add_paragraph()
        p_d.text = s_det
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = COLOR_GRAY_LIGHT
        p_d.space_before = Pt(3)

    # =========================================================================
    # SLIDE 16: Master Process Architecture Conclusion
    # =========================================================================
    slide16 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide16, COLOR_BG_DARK)

    cta_card = slide16.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(1.1), Inches(10.33), Inches(5.3))
    cta_card.fill.solid()
    cta_card.fill.fore_color.rgb = COLOR_CARD_DARK
    cta_card.line.color.rgb = COLOR_RED
    cta_card.line.width = Pt(2)

    tf16 = cta_card.text_frame
    tf16.word_wrap = True
    tf16.margin_left = Inches(0.6)
    tf16.margin_top = Inches(0.5)
    tf16.margin_right = Inches(0.6)

    p = tf16.paragraphs[0]
    p.text = "LIFEQR 2.0 PROCESS ARCHITECTURE MASTERED"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf16.add_paragraph()
    p2.text = "A Complete, Synchronized, Real-World Emergency Medical Network"
    p2.font.size = Pt(17)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_RED
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    p3 = tf16.add_paragraph()
    p3.text = "By unifying:\n✔ Process 1: Zero-battery cryptographic NTAG 424 DNA cards & school student safety tags\n✔ Process 2: Autonomous motorcycle crash IMU detection & paramedic tactical dual-scanning\n✔ Process 3: Real-time in-transit ER telemetry, advance trauma bay reservation & doctor desktop cradle\n... LifeQR transforms emergency survival from a chaotic gamble into an engineered science."
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_GRAY_LIGHT
    p3.alignment = PP_ALIGN.LEFT
    p3.space_before = Pt(18)

    p4 = tf16.add_paragraph()
    p4.text = "READY FOR HARDWARE FABRICATION & MUNICIPAL EMS INTEGRATION"
    p4.font.size = Pt(12)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GREEN
    p4.alignment = PP_ALIGN.CENTER
    p4.space_before = Pt(26)

    # Save presentation
    output_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '../presentation/LifeQR_Process_Architecture_Master.pptx'))
    prs.save(output_path)
    print(f"[Success] Process Architecture Master Presentation generated at: {output_path}")

if __name__ == '__main__':
    create_process_architecture_presentation()
